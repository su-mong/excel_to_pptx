# 라이브러리 추가
from pptx import Presentation
import pandas as pd

# 시작 로그 찍기
print('START')

# 원본이 되는 파일 열기
xlsx = pd.read_excel('origin.xlsx', sheet_name='논문_1992-2021')
prs = Presentation('origin.pptx')

# '논문' 레이아웃에 해당하는 슬라이드 종류를 미리 선택해둔다. (슬라이드 추가 시 논문 레이아웃으로 추가하기 위함)
# 원본 파일을 보면, 3번째 슬라이드 마스터의 9번째 레이아웃이 논문 슬라이드이다.
thesis_slide_layout = prs.slide_masters[2].slide_layouts[8]
# 아래의 반복문 안에서 만들어지는 new_slide를 저장하기 위한 변수.(초기값은 그냥 아무거나 넣었다.)
new_slide = 0

# 엑셀 시트 전체 읽어들이고, 줄별로 반복문 돌리기
for i in range(0, len(xlsx)):
    # 5번째 작업마다 로그를 남김(진행상황 확인용)
    if i % 5 == 0:
        print(str(i) + '번째 작업 진행 중...')

    # 1. 각 줄별로 필요한 데이터를 뽑는다.
    #    각 데이터가 빈칸일 경우 '내용 없음'으로 들어가도록 설정했다.
    data_year = '내용 없음' if xlsx.iloc[i]['연도'] == '' else str(xlsx.iloc[i]['연도'])[0:4]  # 문자열을 자른 이유 : '2000.0'처럼 이상하게 들어가길래 앞의 4자리만 불러오기 위해서 잘랐다.
    data_journal = '내용 없음' if xlsx.iloc[i]['저널명'] == '' else xlsx.iloc[i]['저널명']
    data_title = '내용 없음' if xlsx.iloc[i]['논문제목'] == '' else xlsx.iloc[i]['논문제목']
    data_writers = '내용 없음' if xlsx.iloc[i]['저자'] == '' else xlsx.iloc[i]['저자']
    data_abstract = '내용 없음' if xlsx.iloc[i]['abstract'] == '' else xlsx.iloc[i]['abstract']

    # 2. 논문 레이아웃에 넣을 수 있는 논문 수가 2개이기 때문에, i가 짝수일 때와 홀수일 때의 작업을 나눈다.
    if i % 2 == 0:
        # 3. 논문 레이아웃을 추가한다.
        new_slide = prs.slides.add_slide(thesis_slide_layout)

        # 4. 새 논문 슬라이드의 위쪽에 텍스트를 넣기 위해 placeholder 정보를 불러온다.
        placeholder_year = new_slide.placeholders[19]
        placeholder_journal = new_slide.placeholders[20]
        placeholder_title = new_slide.placeholders[14]
        placeholder_writers = new_slide.placeholders[17]
        placeholder_abstract = new_slide.placeholders[13]

        # 5. 새 논문 슬라이드의 위쪽에 데이터를 넣는다.
        #    (혹시 모를 오류에 대비해서 문자열로 변환했다.)
        placeholder_year.text = str(data_year)
        placeholder_journal.text = str(data_journal)
        placeholder_title.text = str(data_title)
        placeholder_writers.text = str(data_writers)
        placeholder_abstract.text = '내용 없음' if str(data_abstract) == 'nan' else str(data_abstract)

    else:
        # 3. 바로 이전에 추가한 논문 슬라이드의 아래쪽에 텍스트를 넣기 위해 placeholder 정보를 불러온다.
        placeholder_year = new_slide.placeholders[28]
        placeholder_journal = new_slide.placeholders[32]
        placeholder_title = new_slide.placeholders[30]
        placeholder_writers = new_slide.placeholders[26]
        placeholder_abstract = new_slide.placeholders[22]

        # 4. 새 논문 슬라이드의 아래쪽에 데이터를 넣는다.
        #    (혹시 모를 오류에 대비해서 문자열로 변환했다.)
        placeholder_year.text = str(data_year)
        placeholder_journal.text = str(data_journal)
        placeholder_title.text = str(data_title)
        placeholder_writers.text = str(data_writers)
        placeholder_abstract.text = '내용 없음' if str(data_abstract) == 'nan' else str(data_abstract)

# 결과 저장
prs.save('result.pptx')

# 종료 로그 찍기
print('SUCCESS')
